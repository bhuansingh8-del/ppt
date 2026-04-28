import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io
import requests

# --- DETAILED SLIDE CONTENT DATABASE ---
SLIDES = [
    {
        "title": "Digital Activism: Social Action and Social Movements",
        "subtitle": "Course Code 401, M.A. Social Work, Semester IV\nUniversity of Delhi\nPresented to: Dr. Pushpanjali Jha\nGroup Members: Monika R., Aniruddh Singh, Suhail Khan, Shambhavi Mishra, Drisya, Bhawna Sharma, Tannu Tanwar, Sayna Choudhary, Gaurav Nishad, Ankit Kumar",
        "img": "https://images.unsplash.com/photo-1573164713988-8665fc963095?w=800"
    },
    {
        "author": "Monika R.",
        "title": "The Emergence of Digital Activism",
        "content": "• Historical Continuity: Digital activism is a reconfiguration of long-standing practices of collective action rather than a sudden rupture.\n• Technological Evolution: Activism has shifted through print, radio, and television, with each shift redefining the scale and visibility of movements.\n• The Internet’s Dimension: Unlike previous media, the internet enables immediacy, interactivity, and transnational connectivity.\n• Structural Markers: The 1990s Zapatista uprising is a critical historical marker for using early digital networks to mobilize global solidarity.",
        "img": "https://images.unsplash.com/photo-1529156069898-49953e39b3ac?w=800"
    },
    {
        "author": "Aniruddh Singh",
        "title": "Trajectories and Milestones",
        "content": "• 1994 - The Informational Guerrilla: The EZLN (Zapatistas) achieved international reach despite primitive infrastructure.\n• 2004 - The Rise of Web 2.0: The era of 'prosumption' began, where users simultaneously produce and consume activist content.\n• 2011 - The Year of the Protester: Time Magazine designated 2011 as such due to the Arab Spring, Spanish Indignados, and Occupy Wall Street movements.\n• Case Study: Invisible Children utilized YouTube in 2006 to mobilize 80,000 people for the 'Night Commute'.",
        "img": "https://images.unsplash.com/photo-1551818255-e6e10975bc17?w=800"
    },
    {
        "author": "Suhail Khan",
        "title": "Foundational Concepts & Terminologies",
        "content": "• Cyberactivism: The use of internet technologies specifically for progressive political purposes.\n• Hacktivism: A subset of cyberactivism involving technical exploits like disrupting servers for political ends.\n• Electronic Civil Disobedience (ECD): Translates nonviolent rule-breaking into the virtual sphere (e.g., 'virtual sit-ins').\n• Performative Vocabulary: Digital protest acts are often scripted, staged, and enacted for multiple audiences simultaneously.",
        "img": "https://images.unsplash.com/photo-1550751827-4bd374c3f58b?w=800"
    },
    {
        "author": "Suhail Khan",
        "title": "The Choreography of Assembly",
        "content": "• Cyberspace vs. Cyberplace: While cyberspace is abstract, 'Cyberplace' is socially embedded and creates meaningful places for gathering.\n• Choreography of Assembly: The process of using social media to emotionally prepare and symbolically direct physical gatherings.\n• Choreographic Leadership: Indirect, consensual influence exercised by activist elites who 'propose' scripts rather than command.\n• Smart Mobs: Impromptu gatherings organized online, realized physically, and disseminated back into digital spaces.",
        "img": "https://images.unsplash.com/photo-1517048676732-d65bc937f952?w=800"
    },
    {
        "author": "Shambhavi Mishra",
        "title": "Ideological Orientation",
        "content": "• Participatory Democracy: The belief that technology can democratize communication by bypassing corporate gatekeepers.\n• Anti-Corporate Foundations: Protest must move into the digital infrastructures where 'nomadic' corporate power now resides.\n• Horizontalism: The ideal of leaderless, self-organizing networks where no single individual commands the movement.\n• The Public Sphere: Seeking to create online spaces for rational deliberation free from state and market distortions.",
        "img": "https://images.unsplash.com/photo-1540910419892-f0c74b0e896a?w=800"
    },
    {
        "author": "Bhawna Sharma",
        "title": "Global Issues in Focus",
        "content": "• Human Rights: Movements like Black Lives Matter use digital tools to expose localized police brutality to a global audience.\n• Gender Justice: The #MeToo movement empowered survivors to challenge patriarchal norms through visibility.\n• Environmental Justice: Fridays for Future mobilized global youth through digital networks to demand ecological accountability.\n• Digital Sphere Concerns: Activism now includes critiques of data privacy, surveillance, and corporate control over the internet.",
        "img": "https://images.unsplash.com/photo-1611162617213-7d7a39e9b1d7?w=800"
    },
    {
        "author": "Drisya V.",
        "title": "Strategies for Social Change",
        "content": "• Issue-Framing: Using personal stories and visuals to shape narratives and build public empathy.\n• Mobilization: Bridging the online-offline divide to move people from 'liking' a post to attending a physical protest.\n• Advocacy: Directly targeting institutions by flooding inboxes or trending topics to demand policy change.\n• The Quadruple 'A' Approach: Abstaining from platforms, Attacking them, building Alternative tools, or Adapting via clever usage.",
        "img": "https://images.unsplash.com/photo-1460925895917-afdab827c52f?w=800"
    },
    {
        "author": "Drisya V.",
        "title": "Tactics in Practice",
        "content": "• Hashtag Activism: Uniting disparate voices into a global chorus, such as #BlackLivesMatter or #MeToo.\n• Citizen Journalism: Bypassing traditional media during confrontations, such as Facebook Live during Standing Rock.\n• Viral Storytelling: Explaining complex issues through relatable memes, infographics, or short-form videos.\n• Data Activism: Crowdsourcing satellite imagery or citizen reports to expose corporate or environmental wrongdoing.",
        "img": "https://images.unsplash.com/photo-1516321318423-f06f85e504b3?w=800"
    },
    {
        "author": "Ankit Kumar & Sayna Choudhary",
        "title": "Organizational Structures",
        "content": "• Flexible Networks: A shift from rigid hierarchies to fluid, horizontal networks where participants join or leave easily.\n• Key Features: Characterized by decentralization, connectivity across platforms, and an absence of formal bureaucratic rules.\n• Coalitions: Alliances between NGOs, students, and citizens that facilitate resource sharing and wider impact.\n• Technological Structuring: Messaging apps like WhatsApp and Telegram are now central to real-time coordination and planning.",
        "img": "https://images.unsplash.com/photo-1522202176988-66273c2fd55f?w=800"
    },
    {
        "author": "Tannu Tanwar & Gaurav Nishad",
        "title": "Achievements of Digital Activism",
        "content": "• Rapid Mobilization: The ability to organize mass protests and campaigns within hours, as seen during the Arab Spring.\n• Connective Action: Personalized forms of engagement replace hierarchical organizational membership.\n• Creating Counter-Publics: Digital platforms allow marginalized groups to articulate alternative narratives and challenge news hierarchies.\n• Policy Impact: Translating online narratives into offline change, such as new workplace harassment laws resulting from #MeToo.",
        "img": "https://images.unsplash.com/photo-1552664730-d307ca884978?w=800"
    },
    {
        "author": "Tannu Tanwar & Gaurav Nishad",
        "title": "Outcomes and Limitations",
        "content": "• Slacktivism: The risk that low-effort online actions create an illusion of participation without meaningful change.\n• The Digital Divide: Participation is constrained by socioeconomic status, gender, and geographic location.\n• The Double-Edged Sword: The same tools that enable activism also facilitate state surveillance and repression.\n• Sustainability: Digital movements often face rapid cycles of attention where issues are quickly visible but just as quickly forgotten.",
        "img": "https://images.unsplash.com/photo-1504384308090-c894fdcc538d?w=800"
    }
]

def generate_ppt():
    prs = Presentation()
    for item in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content layout
        
        # Format Title
        title_shape = slide.shapes.title
        title_shape.text = item["title"]
        
        # Format Body Text
        body = slide.placeholders[1].text_frame
        body.clear() # Clear default formatting
        
        if "author" in item:
            p = body.add_paragraph()
            p.text = f"Contributor: {item['author']}"
            p.font.bold = True
            p.font.size = Pt(18)
        
        if "content" in item:
            lines = item["content"].split("\n")
            for line in lines:
                p = body.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
                p.space_after = Pt(10)
        elif "subtitle" in item:
            lines = item["subtitle"].split("\n")
            for line in lines:
                p = body.add_paragraph()
                p.text = line
                p.font.size = Pt(18)

        # Add Image from URL (if available)
        try:
            response = requests.get(item["img"], timeout=5)
            if response.status_code == 200:
                img_stream = io.BytesIO(response.content)
                slide.shapes.add_picture(img_stream, Inches(5.5), Inches(2), height=Inches(4))
        except:
            pass # Fails silently if image cannot be loaded

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# --- STREAMLIT UI ---
st.set_page_config(page_title="Digital Activism Tool", page_icon="✊", layout="wide")

st.markdown("<h1 style='text-align: center;'>✊ Digital Activism: Group Presentation Builder</h1>", unsafe_allow_html=True)
st.write("---")

col1, col2, col3 = st.columns([1, 2, 1])
with col2:
    st.info("Click the button below to generate and download the `.pptx` file containing all 12 detailed slides for your assignment.")
    if st.button("🚀 Generate & Download Presentation", use_container_width=True):
        with st.spinner("Compiling assignment data and building slides..."):
            ppt_data = generate_ppt()
            st.success("Presentation generated successfully!")
            st.download_button(
                label="📂 Download PPTX File",
                data=ppt_data,
                file_name="Digital_Activism_Group_Assignment.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )

st.write("---")
st.subheader("Interactive Presentation Preview")
st.write("Click on the tabs below to read the detailed content for each slide. Opening a new tab automatically hides the previous one.")

# Creating Tab Names
tab_names = [
    "1. Intro", "2. Emergence", "3. Trajectories", "4. Concepts", 
    "5. Choreography", "6. Ideology", "7. Global Issues", "8. Strategies", 
    "9. Tactics", "10. Structure", "11. Achievements", "12. Limitations"
]

# Generate Streamlit Tabs
tabs = st.tabs(tab_names)

# Populate Each Tab with Content
for i, tab in enumerate(tabs):
    with tab:
        slide = SLIDES[i]
        
        # Display Image and Text side-by-side inside the tab
        text_col, img_col = st.columns([3, 2])
        
        with text_col:
            st.markdown(f"### {slide['title']}")
            if "author" in slide:
                st.markdown(f"**👤 Contributor:** {slide['author']}")
            st.write("") # Spacer
            
            # Format the content text with markdown blockquotes for clean reading
            if "content" in slide:
                for line in slide['content'].split("\n"):
                    st.write(line)
            else:
                st.write(slide['subtitle'])
                
        with img_col:
            st.image(slide['img'], use_container_width=True)
